import tkinter as tk
from tkinter import filedialog, scrolledtext, messagebox, ttk
import os
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import extract_msg
import fitz
from PIL import Image, ImageTk
import re

opened_files = []  
excel_sheets = {}
current_match_indices = []
current_match_position = 0

def get_current_text_widget():
    if not opened_files:
        return None
    current_tab = notebook.index(notebook.select())
    return opened_files[current_tab]["text_widget"]

def search_and_replace():
    search_term = search_entry.get()
    replace_term = replace_entry.get()
    text_widget = get_current_text_widget()
    if not search_term or not text_widget:
        messagebox.showwarning("Missing input", "Enter a search term.")
        return

    content = text_widget.get(1.0, tk.END)
    pattern = re.escape(search_term)
    matches = re.findall(pattern, content)
    match_count = len(matches)

    if match_count == 0:
        messagebox.showinfo("No Match", f"No occurrences of '{search_term}' found.")
        return

    new_content = re.sub(pattern, replace_term, content)
    text_widget.delete(1.0, tk.END)
    text_widget.insert(tk.END, new_content)
    messagebox.showinfo("Success", f"Replaced {match_count} occurrence(s) of '{search_term}' with '{replace_term}'.")

def find_all_matches():
    global current_match_indices, current_match_position
    text_widget = get_current_text_widget()
    if not text_widget:
        return

    text_widget.tag_remove("highlight", "1.0", tk.END)
    text_widget.tag_remove("active_highlight", "1.0", tk.END)

    search_term = search_entry.get()
    current_match_indices = []
    current_match_position = 0
    if not search_term:
        return

    start = "1.0"
    while True:
        pos = text_widget.search(search_term, start, stopindex=tk.END)
        if not pos:
            break
        end = f"{pos}+{len(search_term)}c"
        current_match_indices.append((pos, end))
        start = end

    for s, e in current_match_indices:
        text_widget.tag_add("highlight", s, e)
    text_widget.tag_config("highlight", background="yellow")

    if current_match_indices:
        focus_current_match()

def focus_current_match():
    text_widget = get_current_text_widget()
    if not text_widget or not current_match_indices:
        return
    text_widget.tag_remove("active_highlight", "1.0", tk.END)
    s, e = current_match_indices[current_match_position]
    text_widget.tag_add("active_highlight", s, e)
    text_widget.tag_config("active_highlight", background="orange")
    text_widget.see(s)

def replace_current():
    global current_match_indices, current_match_position
    text_widget = get_current_text_widget()
    if not text_widget or not current_match_indices:
        return

    s, e = current_match_indices[current_match_position]
    text_widget.delete(s, e)
    text_widget.insert(s, replace_entry.get())
    find_all_matches()
    current_match_position = min(current_match_position, len(current_match_indices) - 1)
    focus_current_match()

def next_match():
    global current_match_position
    if current_match_indices:
        current_match_position = (current_match_position + 1) % len(current_match_indices)
        focus_current_match()

def prev_match():
    global current_match_position
    if current_match_indices:
        current_match_position = (current_match_position - 1) % len(current_match_indices)
        focus_current_match()

def read_docx(file_path):
    doc = Document(file_path)
    return '\n'.join(p.text for p in doc.paragraphs)

def read_xlsx(file_path, sheet_name=None):
    wb = load_workbook(file_path)
    if sheet_name and sheet_name in wb.sheetnames:
        sheets_to_read = [sheet_name]
    else:
        sheets_to_read = wb.sheetnames

    output = []
    for sheet in sheets_to_read:
        ws = wb[sheet]
        output.append(f"--- Sheet: {sheet} ---")
        for row in ws.iter_rows(values_only=True):
            output.append('\t'.join([str(cell) if cell else '' for cell in row]))
    return '\n'.join(output), wb.sheetnames

def read_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return '\n'.join(texts)

def read_text_file(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def read_email(file_path):
    msg = extract_msg.Message(file_path)
    return f"Subject: {msg.subject}\nFrom: {msg.sender}\nTo: {msg.to}\nDate: {msg.date}\n\n{msg.body}"

def read_pdf(file_path):
    text = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text.append(page.get_text())
    return '\n'.join(text)

def read_file(file_path, selected_sheet=None):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.xlsx':
        content, sheets = read_xlsx(file_path, selected_sheet)
        excel_sheets[file_path] = sheets
        return content
    elif ext == '.docx':
        return read_docx(file_path)
    elif ext == '.pptx':
        return read_pptx(file_path)
    elif ext in ['.eml', '.msg']:
        return read_email(file_path)
    elif ext == '.pdf':
        return read_pdf(file_path)
    elif ext in ['.txt', '.c', '.h', '.py', '.js', '.html']:
        return read_text_file(file_path)
    elif ext in ['.png', '.jpg', '.jpeg', '.gif']:
        return "<<<IMAGE_PREVIEW>>>"
    else:
        return f"Unsupported file type: {ext}"

def open_files():
    file_paths = filedialog.askopenfilenames(
        title="Select files",
        filetypes=[("All supported files", "*.docx *.xlsx *.pptx *.txt *.c *.h *.py *.msg *.eml *.pdf")]
    )

    for path in file_paths:
        content = read_file(path)
        tab = tk.Frame(notebook, bg="#fdfefe")

        if content == "<<<IMAGE_PREVIEW>>>":
            img = Image.open(path)
            img = img.resize((600, 400), Image.LANCZOS)
            img_tk = ImageTk.PhotoImage(img)

            label = tk.Label(tab, image=img_tk)
            label.image = img_tk
            label.pack(expand=True)
        else:
            text_area = scrolledtext.ScrolledText(tab, wrap=tk.WORD, font=("Consolas", 10),
                                                  bg="#ffffff", fg="#2c3e50", insertbackground="#2c3e50",
                                                  padx=10, pady=10, bd=0)
            text_area.insert(tk.END, content)
            text_area.pack(fill='both', expand=True)
            opened_files.append({"path": path, "content": content, "text_widget": text_area})

            ext = os.path.splitext(path)[1].lower()
            if ext == '.xlsx' and path in excel_sheets:
                sheet_selector['values'] = excel_sheets[path]
                sheet_selector.set(excel_sheets[path][0])
            else:
                sheet_selector.set('')
                sheet_selector['values'] = []

        notebook.add(tab, text=os.path.basename(path))


def save_pdf(file_path, text):
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text) 
    doc.save(file_path)
    doc.close()

def get_versioned_filename(base_path, ext):
        base_name = os.path.splitext(os.path.basename(base_path))[0]
        directory = os.path.dirname(base_path)
        version = 1
        while True:
            new_name = f"{base_name}_Version{version}{ext}"
            full_path = os.path.join(directory, new_name)
            if not os.path.exists(full_path):
                return full_path
            version += 1 


def save_file():
    if not opened_files:
        messagebox.showwarning("No file", "No file loaded to save.")
        return

    current_tab = notebook.index(notebook.select())
    selected = opened_files[current_tab]
    original_path = selected["path"]
    ext = os.path.splitext(original_path)[1].lower()
    text_widget = selected["text_widget"]
    new_content = text_widget.get(1.0, tk.END).strip()

    filetypes_map = {
        '.pdf': [("PDF files", "*.pdf")],
        '.docx': [("Word Documents", "*.docx")],
        '.xlsx': [("Excel Sheets", "*.xlsx")],
        '.txt': [("Text Files", "*.txt")],
        '.py': [("Python Files", "*.py")],
        '.html': [("HTML Files", "*.html")],
        '.js': [("JavaScript Files", "*.js")],
        '.c': [("C Files", "*.c")],
        '.h': [("Header Files", "*.h")]
    }

    suggested_name = os.path.basename(get_versioned_filename(original_path, ext))
    save_path = filedialog.asksaveasfilename(
        defaultextension=ext,
        initialfile=suggested_name,
        filetypes=filetypes_map.get(ext, [("All Files", "*.*")])
    )

    if not save_path:
        return

    try:
        if ext == ".docx":
            doc = Document()
            for line in new_content.splitlines():
                doc.add_paragraph(line)
            doc.save(save_path)

        elif ext == ".xlsx":
            sheet_name = sheet_selector.get()
            if not sheet_name:
                messagebox.showwarning("Sheet not selected", "Please select a sheet before saving.")
                return

            wb = load_workbook(original_path)
            if sheet_name not in wb.sheetnames:
                messagebox.showerror("Error", f"Sheet '{sheet_name}' not found.")
                return

            ws = wb[sheet_name]
            ws.delete_rows(1, ws.max_row)

            for row_index, line in enumerate(new_content.splitlines(), start=1):
                values = line.split('\t')
                for col_index, value in enumerate(values, start=1):
                    ws.cell(row=row_index, column=col_index, value=value)

            wb.save(save_path)

        elif ext in ['.txt', '.c', '.h', '.py', '.js', '.html']:
            with open(save_path, 'w', encoding='utf-8') as f:
                f.write(new_content)

        elif ext == '.pdf':
            save_pdf(save_path, new_content)

        else:
            messagebox.showinfo("Info", f"Saving for {ext} not supported yet.")
            return

        messagebox.showinfo("Success", f"File saved as: {save_path}")

    except Exception as e:
        messagebox.showerror("Error", f"Could not save file:\n{str(e)}")

def on_sheet_selected(event=None):
    if not opened_files:
        return

    current_tab = notebook.index(notebook.select())
    selected = opened_files[current_tab]["path"]
    sheet_name = sheet_selector.get()
    content = read_file(selected, sheet_name)

    text_widget = opened_files[current_tab]["text_widget"]
    text_widget.delete(1.0, tk.END)
    text_widget.insert(tk.END, content)


def clear_output():
    current_tab = notebook.select()
    if not current_tab:
        return

    index = notebook.index(current_tab)
    notebook.forget(current_tab)
    del opened_files[index]

    
    if not opened_files:
        sheet_selector.set('')
        sheet_selector['values'] = []
    else:
        new_index = min(index, len(opened_files) - 1)
        new_path = opened_files[new_index]["path"]
        ext = os.path.splitext(new_path)[1].lower()
        if ext == ".xlsx" and new_path in excel_sheets:
            sheet_selector['values'] = excel_sheets[new_path]
            sheet_selector.set(excel_sheets[new_path][0])
        else:
            sheet_selector.set('')
            sheet_selector['values'] = []



root = tk.Tk()
root.title("\U0001F4C1 Smart File Reader & Editor Tool")
root.geometry("1000x650")
root.configure(bg="#f0f2f5")

header = tk.Label(root, text="\U0001F4C2 Smart File Reader & Editor", font=("Segoe UI", 22, "bold"), bg="#f0f2f5", fg="#2d3436")
header.pack(pady=15)

button_frame = tk.Frame(root, bg="#f0f2f5")
button_frame.pack(pady=5)

btn = tk.Button(button_frame, text="\U0001F4C2 Open Files", command=open_files,
                font=("Segoe UI", 11, "bold"), bg="#2ecc71", fg="white", width=18)
btn.grid(row=0, column=0, padx=10, pady=6)

save_btn = tk.Button(button_frame, text="\U0001F4BE Save Changes", command=save_file,
                     font=("Segoe UI", 11, "bold"), bg="#3498db", fg="white", width=18)
save_btn.grid(row=0, column=1, padx=10, pady=6)

clear_btn = tk.Button(button_frame, text="\U0001F9F9 Clear Output", command=clear_output,
                      font=("Segoe UI", 11, "bold"), bg="#e67e22", fg="white", width=18)
clear_btn.grid(row=0, column=2, padx=10, pady=6)

search_frame = tk.Frame(root, bg="#f0f2f5")
search_frame.pack(pady=10)

tk.Label(search_frame, text="\U0001F50D Search:", font=("Segoe UI", 10, "bold"), bg="#f0f2f5").grid(row=0, column=0, padx=5)
search_entry = tk.Entry(search_frame, width=35, font=("Segoe UI", 10))
search_entry.grid(row=0, column=1, padx=5)
search_entry.bind("<KeyRelease>", lambda e: find_all_matches())

tk.Label(search_frame, text="\U0001F4DD Replace:", font=("Segoe UI", 10, "bold"), bg="#f0f2f5").grid(row=0, column=2, padx=5)
replace_entry = tk.Entry(search_frame, width=35, font=("Segoe UI", 10))
replace_entry.grid(row=0, column=3, padx=5)

replace_btn = tk.Button(search_frame, text="\U0001F501 Replace All", command=search_and_replace,
                        font=("Segoe UI", 10, "bold"), bg="#9b59b6", fg="white", width=12)
replace_btn.grid(row=0, column=4, padx=8)

navigate_btns = tk.Frame(search_frame, bg="#f0f2f5")
navigate_btns.grid(row=1, column=0, columnspan=5, pady=8)

prev_btn = tk.Button(navigate_btns, text="\u2B05 Prev", command=prev_match,
                     font=("Segoe UI", 10), bg="#7f8c8d", fg="white", width=10)
prev_btn.pack(side="left", padx=5)

replace_one_btn = tk.Button(navigate_btns, text="\u270F Replace This", command=replace_current,
                            font=("Segoe UI", 10), bg="#e67e22", fg="white", width=14)
replace_one_btn.pack(side="left", padx=5)

next_btn = tk.Button(navigate_btns, text="Next \u27A1", command=next_match,
                     font=("Segoe UI", 10), bg="#7f8c8d", fg="white", width=10)
next_btn.pack(side="left", padx=5)

sheet_selector_row = tk.Frame(root, bg="#f0f2f5")
sheet_selector_row.pack(fill='x', padx=20)
sheet_selector_label = tk.Label(sheet_selector_row, text="\U0001F4C4 Select Sheet (Excel):", 
                                font=("Segoe UI", 10, "bold"), bg="#f0f2f5")
sheet_selector_label.pack(side="left", padx=10)

sheet_selector = ttk.Combobox(sheet_selector_row, state="readonly", width=30)
sheet_selector.pack(side="left")
sheet_selector.bind("<<ComboboxSelected>>", on_sheet_selected)


notebook = ttk.Notebook(root)
notebook.pack(fill="both", expand=True, padx=10, pady=10)

footer = tk.Label(
    root,
    text="――――――――――――――――――――――――――\n❤️ Crafted with care by Garima Roy",
    font=("Segoe UI", 9, "italic"),
    bg="#f0f2f5",
    fg="#ff4d4f",  
    justify="center"
)
footer.pack(side="bottom", fill="x", pady=(0, 5))


# Start the application
root.mainloop()
